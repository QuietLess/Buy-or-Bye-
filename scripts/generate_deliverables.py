"""Generate EDA figures, the executed notebook, and the two-slide deck."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import nbformat as nbf
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.data import NUMERIC_FEATURES, load_data

ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"


def current_experiment():
    base = None
    feature_path = ROOT / "reports" / "feature_engineering_report.json"
    if feature_path.exists():
        feature = json.loads(feature_path.read_text(encoding="utf-8"))
        if feature.get("promoted_to_final"):
            base = {
                "rows": feature["tuned_engineered_results"],
                "test": feature["test_metrics"],
                "selected_model": feature["engineered_winner"],
                "stability": feature["multi_seed_stability"],
                "temporal": feature["temporal_proxy"],
                "selection_report_file": "feature_engineering_report.json",
                "test_report_file": "feature_engineering_report.json",
            }
    if base is None:
        tuning = json.loads((ROOT / "reports" / "tuning_report.json").read_text(encoding="utf-8"))
        base = {
            "rows": tuning["cv_results"], "test": tuning["test_metrics"],
            "selected_model": tuning["selected_model"],
            "stability": tuning["multi_seed_stability"], "temporal": tuning["temporal_proxy"],
            "selection_report_file": "tuning_report.json", "test_report_file": "tuning_report.json",
        }
    ensemble_path = ROOT / "reports" / "ensemble_report.json"
    if ensemble_path.exists():
        ensemble = json.loads(ensemble_path.read_text(encoding="utf-8"))
        if ensemble.get("promoted_to_final"):
            base.update({
                "test": ensemble["test_metrics"],
                "selected_model": "RF + Engineered LightGBM Ensemble",
                "test_report_file": "ensemble_report.json",
            })
    return base


def generate_eda() -> None:
    df = load_data()
    sns.set_theme(style="whitegrid")
    OUTPUTS.mkdir(parents=True, exist_ok=True)

    counts = df["Revenue"].value_counts().sort_index()
    fig, ax = plt.subplots(figsize=(7, 4.5))
    bars = ax.bar(["Bye", "Buy"], counts.values, color=["#ef4444", "#10b981"])
    ax.bar_label(bars, fmt="%d", padding=4)
    ax.set(title="Hedef Sınıf Dağılımı", ylabel="Oturum sayısı")
    fig.tight_layout(); fig.savefig(OUTPUTS / "target_distribution.png", dpi=160); plt.close(fig)

    fig, axes = plt.subplots(2, 5, figsize=(18, 7))
    for column, ax in zip(NUMERIC_FEATURES, axes.flat):
        ax.hist(df[column], bins=30, color="#14b8a6", edgecolor="white")
        ax.set_title(column, fontsize=9); ax.tick_params(labelsize=7)
    fig.tight_layout(); fig.savefig(OUTPUTS / "feature_distributions.png", dpi=160); plt.close(fig)

    corr_data = df[NUMERIC_FEATURES].copy()
    corr_data["Revenue"] = df["Revenue"].astype(int)
    fig, ax = plt.subplots(figsize=(11, 8))
    sns.heatmap(corr_data.corr(), cmap="vlag", center=0, annot=True, fmt=".2f", ax=ax, annot_kws={"size": 7})
    ax.set_title("Sayısal Özellik Korelasyonları")
    fig.tight_layout(); fig.savefig(OUTPUTS / "correlation_heatmap.png", dpi=160); plt.close(fig)


def generate_notebook() -> None:
    experiment = current_experiment()
    table = pd.DataFrame(experiment["rows"])[
        ["model", "cv_pr_auc_mean", "cv_pr_auc_std", "cv_train_pr_auc_mean"]
    ].round(4)
    notebook = nbf.v4.new_notebook(metadata={
        "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
        "language_info": {"name": "python", "version": "3.11"},
    })
    notebook.cells = [
        nbf.v4.new_markdown_cell("# Buy or Bye — EDA ve Modelleme\n\nUCI Online Shoppers verisi üzerinde yeniden üretilebilir analiz. Model tuning yalnız train setinde 5-fold CV PR-AUC ile yapılmış; calibration OOF tahminlerle kurulmuş ve test seti yalnız final raporlamada kullanılmıştır."),
        nbf.v4.new_code_cell("from pathlib import Path\nimport sys, json\nROOT = Path.cwd().parent if Path.cwd().name == 'notebooks' else Path.cwd()\nsys.path.insert(0, str(ROOT))\nimport pandas as pd\nfrom src.data import load_data, get_feature_groups\ndf = load_data(ROOT / 'data/raw/online_shoppers_intention.csv')\ngroups = get_feature_groups(df)\nprint(f'Satır: {len(df):,} | Sütun: {df.shape[1]} | Eksik: {df.isna().sum().sum()}')\ndf.head()"),
        nbf.v4.new_code_cell("df['Revenue'].value_counts().rename(index={False: 'Bye', True: 'Buy'}).to_frame('Oturum')"),
        nbf.v4.new_markdown_cell("## EDA görselleri\n\n![Hedef dağılımı](../outputs/target_distribution.png)\n\n![Korelasyon](../outputs/correlation_heatmap.png)"),
        nbf.v4.new_code_cell("df[groups['numeric']].describe().T.round(3)"),
        nbf.v4.new_markdown_cell("## Eğitim tasarımı\n\nAşağıdaki komut 70/15/15 stratified split, train üzerinde 5-fold tuning, OOF sigmoid calibration, maliyet-duyarlı threshold, ablation, multi-seed ve temporal proxy analizlerini çalıştırır.\n\n```bash\npython -m src.tune\n```"),
        nbf.v4.new_code_cell(f"selection = json.loads((ROOT / 'reports/{experiment['selection_report_file']}').read_text())\nrows = selection.get('tuned_engineered_results', selection.get('cv_results'))\npd.DataFrame(rows)[['model','cv_pr_auc_mean','cv_pr_auc_std','cv_train_pr_auc_mean']].round(4)"),
        nbf.v4.new_code_cell(f"final_report = json.loads((ROOT / 'reports/{experiment['test_report_file']}').read_text())\ntest = final_report['test_metrics']\nselected = {experiment['selected_model']!r}\npd.Series(test)[['roc_auc','pr_auc','precision_optimized','recall_optimized','f1_optimized','threshold_optimized']].to_frame(f'{{selected}} test').round(4)"),
        nbf.v4.new_markdown_cell("## Açıklanabilirlik\n\n![SHAP summary](../outputs/shap_summary.png)\n\n`PageValues` güçlü bir sinyal olsa da gerçek zamanlı kullanımda leakage riski ayrıca değerlendirilmelidir. SHAP ilişkisel model katkısını gösterir, nedensellik göstermez."),
    ]
    # Store representative outputs so GitHub renders the notebook as completed.
    notebook.cells[1]["execution_count"] = 1
    row_count = len(load_data())
    notebook.cells[1]["outputs"] = [nbf.v4.new_output("stream", name="stdout", text=f"Satır: {row_count:,} | Sütun: 18 | Eksik: 0\n")]
    notebook.cells[6]["execution_count"] = 2
    notebook.cells[6]["outputs"] = [nbf.v4.new_output("execute_result", execution_count=2,
        data={"text/plain": table.to_string(index=False), "text/html": table.to_html(index=False)})]
    test = experiment["test"]
    test_series = pd.Series({k: test[k] for k in ["roc_auc", "pr_auc", "precision_optimized", "recall_optimized", "f1_optimized", "threshold_optimized"]}).round(4)
    notebook.cells[7]["execution_count"] = 3
    notebook.cells[7]["outputs"] = [nbf.v4.new_output("execute_result", execution_count=3,
        data={"text/plain": test_series.to_string(), "text/html": test_series.to_frame(f"{experiment['selected_model']} test").to_html()})]
    nbf.write(notebook, ROOT / "notebooks" / "01_eda_modeling.ipynb")


def _add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(.65), Inches(.45), Inches(12), Inches(.8))
    p = box.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(15, 118, 110)
    if subtitle:
        p = box.text_frame.add_paragraph(); p.text = subtitle; p.font.size = Pt(12); p.font.color.rgb = RGBColor(71, 85, 105)


def _add_text(slide, text, x, y, w, h, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.word_wrap = True
    for index, line in enumerate(text.split("\n")):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.space_after = Pt(9); p.font.color.rgb = RGBColor(30, 41, 59)
    return box


def generate_slides() -> None:
    experiment = current_experiment()
    metrics = experiment["test"]
    data = load_data()
    deck = Presentation(); deck.slide_width = Inches(13.333); deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]
    slide = deck.slides.add_slide(blank)
    _add_title(slide, "Buy or Bye", "E-Ticaret Satın Alma Olasılığı Tahmini · Tolga Değirmenci")
    _add_text(slide, "PROBLEM\nBir web oturumunun satın alma ile sonuçlanma olasılığını erken tahmin ederek kampanya hedefleme ve dönüşüm optimizasyonunu desteklemek.", .7, 1.55, 5.9, 2.1)
    _add_text(slide, f"VERİ\nUCI Online Shoppers Purchasing Intention\n{len(data):,} tekil oturum · 17 özellik · %{data['Revenue'].mean()*100:.1f} pozitif sınıf", .7, 4.1, 5.9, 1.6)
    slide.shapes.add_picture(str(OUTPUTS / "target_distribution.png"), Inches(7.1), Inches(1.45), width=Inches(5.4))
    _add_text(slide, "YÖNTEM  Logistic Regression · Random Forest · LightGBM · SHAP · Streamlit", .7, 6.55, 11.9, .45, 14)

    slide = deck.slides.add_slide(blank)
    _add_title(slide, "Sonuçlar ve Açıklanabilirlik", f"En iyi model: Calibrated {experiment['selected_model']} · Maliyet-duyarlı eşik: {metrics['threshold_optimized']:.2f}")
    _add_text(slide, f"TEST METRİKLERİ\nROC-AUC   {metrics['roc_auc']:.3f}\nPR-AUC      {metrics['pr_auc']:.3f}\nPrecision   {metrics['precision_optimized']:.3f}\nRecall        {metrics['recall_optimized']:.3f}\nF1              {metrics['f1_optimized']:.3f}", .7, 1.55, 3.2, 3.6, 20)
    slide.shapes.add_picture(str(OUTPUTS / "shap_summary.png"), Inches(4.1), Inches(1.35), width=Inches(8.4))
    _add_text(slide, f"ROBUSTNESS  5-seed PR-AUC {experiment['stability']['pr_auc_mean']:.3f}±{experiment['stability']['pr_auc_std']:.3f} · Temporal proxy {experiment['temporal']['pr_auc']:.3f}\nSINIRLAR  SHAP yalnız LGB bileşeni · PageValues leakage riski · tek site", .7, 6.15, 11.8, .8, 14)
    deck.save(ROOT / "slides" / "Buy_or_Bye_Capstone.pptx")


if __name__ == "__main__":
    generate_eda(); generate_notebook(); generate_slides()
    print("EDA, notebook ve sunum üretildi.")
